from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
DARK_BLUE   = RGBColor(0x00, 0x2D, 0x6E)   # CPPL dark blue
MID_BLUE    = RGBColor(0x00, 0x6E, 0xC7)   # accent
LIGHT_BLUE  = RGBColor(0xE8, 0xF4, 0xFF)   # row fill
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
GREEN       = RGBColor(0x00, 0x7A, 0x4B)
ORANGE      = RGBColor(0xE6, 0x6B, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None):
    """Dark blue top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.3), fill=DARK_BLUE)
    add_text(slide, title,
             Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.82), Inches(11), Inches(0.4),
                 size=14, color=RGBColor(0xB0, 0xCF, 0xFF))
    # thin accent bar
    add_rect(slide, 0, Inches(1.3), W, Inches(0.06), fill=MID_BLUE)


def bullet_block(slide, x, y, w, items, title=None, title_color=DARK_BLUE):
    """Renders a titled bullet list; returns bottom y."""
    cur_y = y
    if title:
        add_text(slide, title, x, cur_y, w, Inches(0.4),
                 size=13, bold=True, color=title_color)
        cur_y += Inches(0.38)
    for item in items:
        add_text(slide, f"  •  {item}", x, cur_y, w, Inches(0.34),
                 size=12, color=DARK_GRAY)
        cur_y += Inches(0.32)
    return cur_y


def tag(slide, text, x, y, fill=MID_BLUE, text_color=WHITE, size=11):
    w = Inches(len(text) * 0.11 + 0.3)
    h = Inches(0.3)
    r = add_rect(slide, x, y, w, h, fill=fill)
    r.line.fill.background()
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return w


def info_card(slide, x, y, w, h, heading, lines, head_fill=MID_BLUE):
    add_rect(slide, x, y, w, Inches(0.4), fill=head_fill)
    add_text(slide, heading, x + Inches(0.1), y + Inches(0.04),
             w - Inches(0.2), Inches(0.34),
             size=12, bold=True, color=WHITE)
    add_rect(slide, x, y + Inches(0.4), w, h - Inches(0.4),
             fill=LIGHT_BLUE, line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width=Pt(1))
    cy = y + Inches(0.48)
    for line in lines:
        add_text(slide, f"  •  {line}", x + Inches(0.05), cy,
                 w - Inches(0.1), Inches(0.32), size=11, color=DARK_GRAY)
        cy += Inches(0.31)


def table_slide(slide, x, y, w, col_widths, headers, rows,
                head_fill=DARK_BLUE, alt_fill=LIGHT_BLUE):
    row_h = Inches(0.38)
    # header row
    cx = x
    for i, hdr in enumerate(headers):
        add_rect(slide, cx, y, col_widths[i], row_h, fill=head_fill)
        add_text(slide, hdr, cx + Inches(0.08), y + Inches(0.06),
                 col_widths[i] - Inches(0.1), row_h - Inches(0.1),
                 size=11, bold=True, color=WHITE)
        cx += col_widths[i]
    # data rows
    for ri, row in enumerate(rows):
        ry = y + row_h * (ri + 1)
        fill = alt_fill if ri % 2 == 0 else WHITE
        cx = x
        for ci, cell in enumerate(row):
            add_rect(slide, cx, ry, col_widths[ci], row_h,
                     fill=fill,
                     line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width=Pt(0.5))
            add_text(slide, str(cell),
                     cx + Inches(0.08), ry + Inches(0.06),
                     col_widths[ci] - Inches(0.1), row_h - Inches(0.08),
                     size=11, color=DARK_GRAY)
            cx += col_widths[ci]


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# full background
add_rect(s, 0, 0, W, H, fill=DARK_BLUE)
# decorative right panel
add_rect(s, Inches(8.8), 0, Inches(4.53), H, fill=MID_BLUE)
# diagonal accent
add_rect(s, Inches(8.2), 0, Inches(0.8), H, fill=RGBColor(0x00, 0x55, 0xAA))

add_text(s, "CPPL", Inches(0.6), Inches(1.4), Inches(7.5), Inches(1.1),
         size=54, bold=True, color=WHITE)
add_text(s, "Inbound Process Automation",
         Inches(0.6), Inches(2.4), Inches(7.8), Inches(0.7),
         size=26, bold=False, color=RGBColor(0xB0, 0xCF, 0xFF))
add_text(s, "Prototype Planning — Data Requirements",
         Inches(0.6), Inches(3.05), Inches(7.8), Inches(0.55),
         size=20, color=RGBColor(0x90, 0xB8, 0xFF))

# divider
add_rect(s, Inches(0.6), Inches(3.7), Inches(4), Inches(0.05), fill=MID_BLUE)

add_text(s, "Webteam Private Limited  ×  TransExcel Consulting",
         Inches(0.6), Inches(3.85), Inches(7.5), Inches(0.4),
         size=13, color=RGBColor(0x80, 0xA8, 0xE0))
add_text(s, "May 2026",
         Inches(0.6), Inches(4.25), Inches(4), Inches(0.35),
         size=12, color=RGBColor(0x70, 0x98, 0xCC))

# right panel content
add_text(s, "Modules in Scope",
         Inches(9.1), Inches(1.8), Inches(3.8), Inches(0.45),
         size=14, bold=True, color=WHITE)
for i, m in enumerate(["M01 — Unique Item Code Creation", "M02 — Material Inward & GRN"]):
    add_text(s, m, Inches(9.1), Inches(2.35 + i * 0.55), Inches(3.9), Inches(0.45),
             size=13, color=WHITE)

add_text(s, "Integration",
         Inches(9.1), Inches(3.6), Inches(3.8), Inches(0.4),
         size=14, bold=True, color=WHITE)
add_text(s, "Tally Prime  →  Our System  →  Tally Prime",
         Inches(9.1), Inches(4.05), Inches(3.9), Inches(0.45),
         size=12, color=RGBColor(0xCC, 0xE5, 0xFF))

add_text(s, "Confidential — Pre-Sales",
         Inches(0.5), Inches(7.05), Inches(5), Inches(0.3),
         size=10, color=RGBColor(0x60, 0x88, 0xBB), italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Agenda", "What we will cover in this session")

items = [
    ("01", "What We Are Building",        "Overview of 2 modules and their purpose"),
    ("02", "How Tally Prime Connects",     "Data flow between Tally and our system"),
    ("03", "Data Needed — Item Code",      "What we need from CPPL for Module 1"),
    ("04", "Data Needed — GRN",           "What we need from CPPL for Module 2"),
    ("05", "Minimum to Start Prototype",  "5 things to kick off immediately"),
    ("06", "Next Steps",                  "Actions, owners, timeline"),
]

for i, (num, title, sub) in enumerate(items):
    row = i // 2
    col = i %  2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.7 + row * 1.6)
    add_rect(s, x, y, Inches(6.0), Inches(1.35),
             fill=LIGHT_BLUE,
             line_color=RGBColor(0xBB, 0xD6, 0xF0), line_width=Pt(1))
    add_rect(s, x, y, Inches(0.65), Inches(1.35), fill=MID_BLUE)
    add_text(s, num, x + Inches(0.05), y + Inches(0.35),
             Inches(0.55), Inches(0.6),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x + Inches(0.75), y + Inches(0.15),
             Inches(5.1), Inches(0.45), size=14, bold=True, color=DARK_BLUE)
    add_text(s, sub, x + Inches(0.75), y + Inches(0.6),
             Inches(5.1), Inches(0.55), size=11, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What We Are Building
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "What We Are Building", "Phase I — 2 modules for client prototype")

# M01 card
add_rect(s, Inches(0.4), Inches(1.6), Inches(5.9), Inches(5.3),
         fill=WHITE, line_color=RGBColor(0xBB, 0xD6, 0xF0), line_width=Pt(1.5))
add_rect(s, Inches(0.4), Inches(1.6), Inches(5.9), Inches(0.55), fill=DARK_BLUE)
add_text(s, "M01   Unique Item Code Creation",
         Inches(0.55), Inches(1.65), Inches(5.6), Inches(0.45),
         size=14, bold=True, color=WHITE)

m01_points = [
    "Hierarchical classification: Type → Category → Class → Characteristics",
    "Distinguishes Service vs Goods; RM vs SFG vs FG",
    "Alphanumeric code — fixed length, intelligently structured",
    "Deduplication algorithm — prevents duplicate items",
    "Description field separate from code (human-readable)",
    "Syncs new items back to Tally as Stock Items",
]
cy = Inches(2.35)
for pt in m01_points:
    add_text(s, f"  •  {pt}", Inches(0.55), cy, Inches(5.6), Inches(0.35),
             size=11, color=DARK_GRAY)
    cy += Inches(0.34)

tag(s, "Tally → Pull Stock Items & Groups", Inches(0.55), Inches(6.55),
    fill=GREEN, size=10)
tag(s, "Push new items back to Tally", Inches(0.55), Inches(6.92),
    fill=ORANGE, size=10)

# M02 card
add_rect(s, Inches(7.0), Inches(1.6), Inches(5.9), Inches(5.3),
         fill=WHITE, line_color=RGBColor(0xBB, 0xD6, 0xF0), line_width=Pt(1.5))
add_rect(s, Inches(7.0), Inches(1.6), Inches(5.9), Inches(0.55), fill=MID_BLUE)
add_text(s, "M02   Material Inward & GRN",
         Inches(7.15), Inches(1.65), Inches(5.6), Inches(0.45),
         size=14, bold=True, color=WHITE)

m02_points = [
    "Receive material against open Purchase Orders from Tally",
    "Auto-fill GRN lines from PO (item, qty, rate, UOM)",
    "QC / Inspection checklist — pass or fail per line item",
    "Split Accepted qty vs Rejected qty",
    "Partial delivery tracking — pending qty remains open",
    "Attach documents — Challan, Invoice, E-way Bill",
    "After approval → push Receipt Note back to Tally",
]
cy = Inches(2.35)
for pt in m02_points:
    add_text(s, f"  •  {pt}", Inches(7.15), cy, Inches(5.6), Inches(0.35),
             size=11, color=DARK_GRAY)
    cy += Inches(0.34)

tag(s, "Tally → Pull POs, Suppliers, Godowns", Inches(7.15), Inches(6.55),
    fill=GREEN, size=10)
tag(s, "Push Receipt Note back to Tally", Inches(7.15), Inches(6.92),
    fill=ORANGE, size=10)

# divider
add_rect(s, Inches(6.5), Inches(1.6), Inches(0.05), Inches(5.3), fill=RGBColor(0xCC, 0xDD, 0xEE))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How Tally Prime Connects
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "How Tally Prime Connects", "Data flow between Tally Prime and our system")

# 3 columns: Pull | Our System | Push
cols = [
    (Inches(0.4),  "PULL from Tally",       DARK_BLUE, GREEN,
     ["Supplier Ledgers\n(Sundry Creditors)",
      "Open Purchase Orders\n& PO line items",
      "Stock Items\n& Stock Groups",
      "Godowns (Warehouses)",
      "Units of Measure",
      "Batches / Lots\n(if enabled)"]),
    (Inches(4.75), "OUR SYSTEM ADDS",        MID_BLUE, MID_BLUE,
     ["Approval workflow\n(Draft → QC → Submit)",
      "QC inspection checklist\nper line item",
      "Accepted vs Rejected\nqty split",
      "Document attachments\n(Challan, Invoice, E-way)",
      "GRN reference number\n(our own series)",
      "Real-time dashboard\n& pending PO tracker"]),
    (Inches(9.1),  "PUSH back to Tally",    DARK_BLUE, ORANGE,
     ["Receipt Note voucher\n(inventory updated in Tally)",
      "Stock Item creation\n(from M01 item codes)",
      "Batch / Lot entries\n(if applicable)",
      "Narration & reference\nnumbers carried over",
      "Triggered only after\nGRN approval — not on save"]),
]

arrows_x = [Inches(4.35), Inches(8.7)]
for ax in arrows_x:
    add_rect(s, ax, Inches(3.8), Inches(0.4), Inches(0.06), fill=MID_BLUE)
    # arrowhead (small triangle approximated by narrow rect)
    add_text(s, "▶", ax + Inches(0.28), Inches(3.68),
             Inches(0.3), Inches(0.3), size=16, color=MID_BLUE)

for x, heading, hfill, badge_fill, points in cols:
    w = Inches(4.0)
    add_rect(s, x, Inches(1.55), w, Inches(0.5), fill=hfill)
    add_text(s, heading, x + Inches(0.1), Inches(1.6),
             w - Inches(0.15), Inches(0.42),
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, Inches(2.05), w, Inches(5.0),
             fill=LIGHT_BLUE,
             line_color=RGBColor(0xBB, 0xD6, 0xF0), line_width=Pt(1))
    cy = Inches(2.2)
    for pt in points:
        add_rect(s, x + Inches(0.1), cy, w - Inches(0.2), Inches(0.62),
                 fill=WHITE,
                 line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width=Pt(0.5))
        add_text(s, pt, x + Inches(0.2), cy + Inches(0.05),
                 w - Inches(0.35), Inches(0.56),
                 size=10, color=DARK_GRAY)
        cy += Inches(0.72)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Data Required: Item Code Module
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Data Required — Item Code Module (M01)",
             "What we need from CPPL to build and demo this module")

# Left: from Tally
info_card(s, Inches(0.4), Inches(1.6), Inches(5.9), Inches(2.8),
          "From Tally Prime  —  Export as Excel / CSV",
          ["Stock Items list  (all existing materials)",
           "Stock Groups / Categories hierarchy",
           "Units of Measure (UOM) list",
           "Any existing item codes or numbering pattern"],
          head_fill=GREEN)

# Right: from CPPL team
info_card(s, Inches(7.0), Inches(1.6), Inches(5.9), Inches(2.8),
          "From CPPL Team  —  Discussion / Manual Input",
          ["How do you currently classify materials?",
           "Do you distinguish RM / SFG / FG today?",
           "5–10 sample items used daily (real names)",
           "Do you track Batch / Heat numbers?"],
          head_fill=DARK_BLUE)

# Bottom: why it matters
add_rect(s, Inches(0.4), Inches(4.65), Inches(12.55), Inches(0.45), fill=DARK_BLUE)
add_text(s, "Why this matters",
         Inches(0.55), Inches(4.69), Inches(4), Inches(0.38),
         size=13, bold=True, color=WHITE)

reasons = [
    ("Prevents duplicates", "The algorithm checks against existing Tally items — we need the full list."),
    ("Realistic demo",      "Real material names from CPPL make the prototype credible to management."),
    ("Hierarchy design",    "We build the code structure around how CPPL actually classifies materials."),
]
for i, (title, desc) in enumerate(reasons):
    x = Inches(0.4 + i * 4.22)
    add_rect(s, x, Inches(5.1), Inches(4.0), Inches(1.9),
             fill=LIGHT_BLUE,
             line_color=RGBColor(0xBB, 0xD6, 0xF0), line_width=Pt(1))
    add_text(s, title, x + Inches(0.15), Inches(5.2),
             Inches(3.7), Inches(0.38), size=12, bold=True, color=DARK_BLUE)
    add_text(s, desc, x + Inches(0.15), Inches(5.6),
             Inches(3.7), Inches(1.1), size=11, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Required: GRN Module
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Data Required — GRN Module (M02)",
             "What we need from CPPL to build and demo this module")

info_card(s, Inches(0.4), Inches(1.6), Inches(5.9), Inches(3.2),
          "From Tally Prime  —  Export as Excel / CSV",
          ["Supplier Ledgers  (Sundry Creditors group)",
           "2–3 open Purchase Orders with line items",
           "Godowns list  (warehouse / store locations)",
           "Units of Measure",
           "Batch / Lot names  (if batch tracking is on)"],
          head_fill=GREEN)

info_card(s, Inches(7.0), Inches(1.6), Inches(5.9), Inches(3.2),
          "From CPPL Team  —  Discussion / Manual Input",
          ["Who approves a GRN?  (role / person name)",
           "What QC checks are done on inbound material?",
           "Rejection handling — return to vendor or hold in store?",
           "Sample challan or delivery document  (any format)",
           "Do partial deliveries happen frequently?"],
          head_fill=DARK_BLUE)

# workflow preview strip
add_rect(s, Inches(0.4), Inches(5.05), Inches(12.55), Inches(0.42), fill=MID_BLUE)
add_text(s, "Proposed GRN Workflow",
         Inches(0.55), Inches(5.09), Inches(4), Inches(0.34),
         size=13, bold=True, color=WHITE)

steps = ["Create GRN\n(Draft)", "Attach\nDocuments", "QC\nInspection",
         "Submit\nfor Approval", "Approved", "Push Receipt\nNote to Tally"]
fills = [LIGHT_BLUE, LIGHT_BLUE, LIGHT_BLUE, LIGHT_BLUE,
         RGBColor(0xD4, 0xED, 0xDA), RGBColor(0xFF, 0xE8, 0xCC)]
bcolors = [MID_BLUE]*4 + [GREEN, ORANGE]

for i, (step, fill, bc) in enumerate(zip(steps, fills, bcolors)):
    x = Inches(0.4 + i * 2.1)
    add_rect(s, x, Inches(5.7), Inches(1.85), Inches(1.4),
             fill=fill, line_color=bc, line_width=Pt(1.5))
    add_text(s, step, x + Inches(0.08), Inches(5.85),
             Inches(1.7), Inches(1.05),
             size=11, bold=(i >= 4), color=bc, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, "→", Inches(0.4 + i * 2.1 + 1.82), Inches(6.1),
                 Inches(0.28), Inches(0.4),
                 size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Minimum to Start
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Minimum to Start Prototype",
             "These 5 inputs are enough for us to begin building immediately")

items = [
    (MID_BLUE,  "01", "Stock Items + Stock Groups export from Tally",
     "Excel / CSV  •  Needed for Item Code module seed data and deduplication"),
    (MID_BLUE,  "02", "Supplier Ledgers export from Tally",
     "Excel / CSV  •  Sundry Creditors group only  •  Needed for GRN supplier dropdown"),
    (MID_BLUE,  "03", "2–3 sample Purchase Orders",
     "Even dummy POs in CPPL format  •  To demo the GRN 'receive against PO' flow"),
    (GREEN,     "04", "5–10 real material names CPPL uses daily",
     "Verbal or written  •  Makes the prototype demo credible to CPPL management"),
    (ORANGE,    "05", "Name of person / role who approves GRN",
     "One name is enough  •  We wire the approval workflow around this role"),
]

for i, (clr, num, title, sub) in enumerate(items):
    y = Inches(1.65 + i * 1.05)
    add_rect(s, Inches(0.4), y, Inches(12.55), Inches(0.92),
             fill=WHITE,
             line_color=RGBColor(0xCC, 0xDD, 0xEE), line_width=Pt(1))
    add_rect(s, Inches(0.4), y, Inches(0.7), Inches(0.92), fill=clr)
    add_text(s, num, Inches(0.4), y + Inches(0.22),
             Inches(0.7), Inches(0.5),
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.25), y + Inches(0.1),
             Inches(11.4), Inches(0.38), size=14, bold=True, color=DARK_BLUE)
    add_text(s, sub, Inches(1.25), y + Inches(0.5),
             Inches(11.4), Inches(0.35), size=11, color=MID_GRAY)

add_rect(s, Inches(0.4), Inches(7.0), Inches(12.55), Inches(0.06), fill=MID_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Next Steps
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Next Steps", "Actions and owners after this call")

table_slide(
    s,
    x=Inches(0.4), y=Inches(1.65), w=Inches(12.55),
    col_widths=[Inches(0.6), Inches(5.5), Inches(3.2), Inches(3.25)],
    headers=["#", "Action", "Owner", "By When"],
    rows=[
        ["1", "Export Stock Items + Stock Groups from Tally",           "CPPL — Accounts / Store",    "Before next call"],
        ["2", "Export Supplier Ledgers (Sundry Creditors) from Tally",  "CPPL — Accounts",            "Before next call"],
        ["3", "Share 2–3 sample POs (or dummy in same format)",         "CPPL — Purchase Team",       "Before next call"],
        ["4", "Confirm GRN approver role / name",                       "CPPL — Management",          "On this call"],
        ["5", "Share 5–10 daily-use material names",                    "CPPL — Store Team",          "On this call"],
        ["6", "Confirm Batch / Lot tracking status in Tally",           "CPPL — Accounts",            "On this call"],
        ["7", "Set up Tally Prime API access / port for Webteam",       "CPPL IT + Webteam",          "After data received"],
        ["8", "Webteam: scaffold project + integrate Tally sync",       "Webteam",                    "Within 1 week"],
        ["9", "Webteam: prototype demo — both modules",                 "Webteam",                    "Within 3 weeks"],
    ]
)

add_rect(s, Inches(0.4), Inches(6.55), Inches(12.55), Inches(0.65),
         fill=DARK_BLUE)
add_text(s, "We are ready to start the moment the data arrives.  "
            "The prototype will be a fully working demo — not a mockup.",
         Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.5),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=DARK_BLUE)
add_rect(s, 0, 0, Inches(4.5), H, fill=MID_BLUE)
add_rect(s, Inches(4.5), 0, Inches(0.08), H, fill=RGBColor(0x00, 0x55, 0xAA))

add_text(s, "Thank You",
         Inches(0.5), Inches(2.2), Inches(3.7), Inches(1.1),
         size=42, bold=True, color=WHITE)
add_rect(s, Inches(0.5), Inches(3.3), Inches(2.5), Inches(0.06), fill=WHITE)
add_text(s, "We look forward to building\nthis with CPPL.",
         Inches(0.5), Inches(3.5), Inches(3.6), Inches(1.0),
         size=14, color=RGBColor(0xCC, 0xE5, 0xFF))

add_text(s, "Webteam Private Limited",
         Inches(5.2), Inches(2.5), Inches(7.5), Inches(0.55),
         size=22, bold=True, color=WHITE)
add_text(s, "TransExcel Consulting",
         Inches(5.2), Inches(3.05), Inches(7.5), Inches(0.45),
         size=18, color=RGBColor(0xB0, 0xCF, 0xFF))

add_rect(s, Inches(5.2), Inches(3.65), Inches(6.5), Inches(0.04),
         fill=RGBColor(0x44, 0x77, 0xAA))

contacts = [
    ("Email",   "arvind.warule@webteam.in"),
    ("Project", "CPPL Inbound Process Automation — Phase I"),
    ("Repo",    "github.com/AviWarule/CPPL"),
]
for i, (label, val) in enumerate(contacts):
    add_text(s, label, Inches(5.2), Inches(3.85 + i * 0.55),
             Inches(1.4), Inches(0.42), size=11,
             color=RGBColor(0x80, 0xA8, 0xD8), bold=True)
    add_text(s, val, Inches(6.5), Inches(3.85 + i * 0.55),
             Inches(6.5), Inches(0.42), size=12, color=WHITE)

add_text(s, "Confidential — Pre-Sales Document",
         Inches(5.2), Inches(6.9), Inches(7.5), Inches(0.35),
         size=10, color=RGBColor(0x55, 0x80, 0xAA), italic=True)


# ── save ────────────────────────────────────────────────────────────────────
out = r"C:\Users\Dell\CPPL\CPPL_Prototype_DataRequirements.pptx"
prs.save(out)
print(f"Saved: {out}")
